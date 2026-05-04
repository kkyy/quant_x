"""
生成工商经营客群授信框架演示文稿（美化版）
输出：/tmp/credit_framework.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── 颜色系统 ──────────────────────────────────────────────────────────────────
DEEP_BLUE    = RGBColor(0x1A, 0x3A, 0x6B)   # 主深蓝
MID_BLUE     = RGBColor(0x25, 0x5F, 0xAB)   # 中蓝（标题背景）
ACCENT_BLUE  = RGBColor(0x2E, 0x86, 0xC1)   # 强调蓝
LIGHT_BLUE   = RGBColor(0xD6, 0xEA, 0xF8)   # 浅蓝背景
ORANGE       = RGBColor(0xE8, 0x7E, 0x04)   # 橙色高亮
GREEN        = RGBColor(0x1A, 0x8A, 0x5A)   # 绿色（√）
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x1C, 0x1C, 0x1C)
MID_TEXT     = RGBColor(0x44, 0x44, 0x55)
CARD_BG1     = RGBColor(0xEB, 0xF5, 0xFB)   # 企业基本面卡片
CARD_BG2     = RGBColor(0xE8, 0xF8, 0xF0)   # 个人基本面卡片
CARD_BG3     = RGBColor(0xFE, 0xF9, 0xE7)   # 授信敏感度卡片
DIVIDER      = RGBColor(0xCC, 0xD9, 0xEA)
BADGE_BG     = RGBColor(0xF0, 0xF4, 0xFC)


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None,
                     line_width=Pt(1), corner=0.1):
    from pptx.oxml.ns import qn
    from lxml import etree
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    # 设置圆角弧度
    sp = shape.element
    prstGeom = sp.spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is not None:
            for gd in avLst.findall(qn('a:gd')):
                if gd.get('name') == 'adj':
                    gd.set('fmla', f'val {int(corner * 50000)}')

    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def set_tf(shape, text, size, bold=False, color=None, align=PP_ALIGN.LEFT,
           word_wrap=True, space_before=0, space_after=0):
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color else DARK_TEXT
    pPr = p._pPr
    if pPr is None:
        from lxml import etree
        from pptx.oxml.ns import qn
        pPr = etree.SubElement(p._p, qn('a:pPr'))
    from pptx.oxml.ns import qn
    pPr.set('spcBef', str(int(space_before * 12700)))
    pPr.set('spcAft', str(int(space_after * 12700)))
    return tf


def add_text_box(slide, l, t, w, h, text, size=11, bold=False,
                 color=None, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    set_tf(txBox, text, size, bold=bold, color=color, align=align,
           word_wrap=word_wrap)
    return txBox


def add_multiline_textbox(slide, l, t, w, h, lines, sizes=None, bolds=None,
                          colors=None, aligns=None, line_spacing_pt=2):
    """多行文本框，lines 是 list of str"""
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        sz = (sizes[i] if sizes else 11)
        bd = (bolds[i] if bolds else False)
        cl = (colors[i] if colors else DARK_TEXT)
        al = (aligns[i] if aligns else PP_ALIGN.LEFT)

        p.alignment = al
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = cl

        # 行间距
        pPr = p._pPr
        if pPr is None:
            pPr = etree.SubElement(p._p, qn('a:pPr'))
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(sz * 100 + line_spacing_pt * 100)))
    return txBox


def add_connector(slide, x1, y1, x2, y2, color=None, width=Pt(1.5)):
    """添加直线连接符"""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line = connector.line
    line.color.rgb = color if color else ACCENT_BLUE
    line.width = width
    return connector


# ── 幻灯片 1：封面 ─────────────────────────────────────────────────────────────
def make_cover(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    W, H = 13.33, 7.5

    # 深蓝底部横条（占下1/3）
    add_rect(slide, 0, H * 0.62, W, H * 0.38, fill_rgb=DEEP_BLUE)
    # 主蓝顶部横条（占上2/3）
    add_rect(slide, 0, 0, W, H * 0.62, fill_rgb=MID_BLUE)
    # 装饰斜线条（右上角）
    for i in range(5):
        add_rect(slide, W - 1.5 + i * 0.25, 0, 0.12, H * 0.62,
                 fill_rgb=RGBColor(0x3A, 0x7A, 0xC8))

    # 橙色左侧竖条
    add_rect(slide, 0.5, H * 0.18, 0.07, H * 0.44, fill_rgb=ORANGE)

    # 主标题
    add_multiline_textbox(
        slide, 0.8, H * 0.2, W - 1.5, 1.4,
        ["工商经营客群", "授信框架优化方案"],
        sizes=[36, 36], bolds=[True, True],
        colors=[WHITE, WHITE],
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT],
        line_spacing_pt=4
    )

    # 副标题
    add_text_box(slide, 0.8, H * 0.52, W - 2, 0.4,
                 "强化经营特质挖掘  ·  引入授信敏感度调整  ·  提升风控精准度",
                 size=13, color=RGBColor(0xAE, 0xD6, 0xF1))

    # 底部说明
    add_text_box(slide, 0.8, H * 0.68, 6, 0.35,
                 "信贷策略研究  |  2026",
                 size=11, color=RGBColor(0x85, 0xA8, 0xD0))

    return slide


# ── 幻灯片 2：客群背景与痛点 ─────────────────────────────────────────────────
def make_background(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    W, H = 13.33, 7.5

    # 顶部蓝色标题栏
    add_rect(slide, 0, 0, W, 1.1, fill_rgb=DEEP_BLUE)
    add_rect(slide, 0, 1.1, W, 0.06, fill_rgb=ORANGE)
    add_text_box(slide, 0.4, 0.22, W - 0.8, 0.65,
                 "01  客群背景与痛点", size=22, bold=True, color=WHITE)
    add_text_box(slide, 0.4, 0.7, W - 0.8, 0.35,
                 "工商经营客群的核心特征与信贷挑战", size=12, color=RGBColor(0xAE, 0xD6, 0xF1))

    # 左侧：客群定位框
    add_rounded_rect(slide, 0.3, 1.4, 3.2, 5.7, fill_rgb=LIGHT_BLUE,
                     line_rgb=ACCENT_BLUE, line_width=Pt(1.5), corner=0.08)
    add_rect(slide, 0.3, 1.4, 3.2, 0.6, fill_rgb=ACCENT_BLUE)
    add_text_box(slide, 0.5, 1.48, 2.8, 0.45,
                 "客群定位", size=13, bold=True, color=WHITE)

    char_items = [
        ("工商经营属性", "具有明确的工商注册记录\n经营年限 ≥ 1 年"),
        ("收入特点", "以经营性收入为主\n无稳定工资流水"),
        ("风险水平", "年化不良率  4.3%\n高于普通零售贷款"),
    ]
    for idx, (title, desc) in enumerate(char_items):
        y = 2.2 + idx * 1.55
        add_rounded_rect(slide, 0.5, y, 2.8, 1.3,
                         fill_rgb=WHITE, line_rgb=DIVIDER, line_width=Pt(1), corner=0.08)
        add_text_box(slide, 0.65, y + 0.08, 2.5, 0.35,
                     title, size=11, bold=True, color=MID_BLUE)
        add_text_box(slide, 0.65, y + 0.4, 2.5, 0.7,
                     desc, size=10, color=MID_TEXT)

    # 右侧三大痛点
    pain_data = [
        (CARD_BG1, ACCENT_BLUE, "⚠ 痛点一：收入难以核实",
         "工商经营客户缺乏稳定工资数据，\n传统收入核验方式失效，\n授信额度评估误差大。"),
        (CARD_BG3, ORANGE, "⚠ 痛点二：风险识别粗糙",
         "行业、经营规模、对外投资状况\n等经营维度未被有效挖掘，\n导致风险区分度不足。"),
        (CARD_BG2, GREEN, "⚠ 痛点三：敏感度未分层",
         "不同客户对授信变动的反应差异\n显著，统一策略造成优质客户流失\n和风险客户留存并存。"),
    ]
    for idx, (bg, accent, title, body) in enumerate(pain_data):
        x, y = 3.85, 1.4 + idx * 1.95
        add_rounded_rect(slide, x, y, 9.1, 1.75,
                         fill_rgb=bg, line_rgb=accent, line_width=Pt(1.5), corner=0.08)
        add_rect(slide, x, y, 0.07, 1.75, fill_rgb=accent)
        add_text_box(slide, x + 0.2, y + 0.12, 8.7, 0.4,
                     title, size=12, bold=True, color=accent)
        add_text_box(slide, x + 0.2, y + 0.52, 8.5, 1.0,
                     body, size=10.5, color=MID_TEXT)

    return slide


# ── 幻灯片 3：整体框架（核心页） ──────────────────────────────────────────────
def make_framework(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    W, H = 13.33, 7.5

    # 顶部
    add_rect(slide, 0, 0, W, 1.1, fill_rgb=DEEP_BLUE)
    add_rect(slide, 0, 1.1, W, 0.06, fill_rgb=ORANGE)
    add_text_box(slide, 0.4, 0.22, W - 0.8, 0.65,
                 "02  客群授信整体框架", size=22, bold=True, color=WHITE)
    add_text_box(slide, 0.4, 0.7, W - 0.8, 0.35,
                 "三维评估体系 — 企业基本面 · 个人基本面 · 授信敏感度", size=12,
                 color=RGBColor(0xAE, 0xD6, 0xF1))

    # 中央标题节点
    add_rounded_rect(slide, 4.8, 1.45, 3.7, 0.85,
                     fill_rgb=MID_BLUE, line_rgb=DEEP_BLUE, line_width=Pt(2), corner=0.12)
    add_text_box(slide, 4.85, 1.52, 3.6, 0.7,
                 "工商经营授信客群", size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # 三个分支卡片 ── 数据
    cards = [
        {
            "bg": CARD_BG1, "accent": ACCENT_BLUE,
            "icon": "🏢", "title": "企业基本面",
            "subtitle": "经营合规性 & 规模评估",
            "items": [
                "✦ 工商特许经营资质核验",
                "✦ 分支机构数量及地域分布",
                "✦ 对外投资与关联企业图谱",
                "✦ 经营年限与续营稳定性",
            ],
            "x": 0.3, "y": 2.65, "w": 3.8, "h": 4.2,
        },
        {
            "bg": CARD_BG2, "accent": GREEN,
            "icon": "👤", "title": "个人基本面",
            "subtitle": "资产负债 & 还款能力",
            "items": [
                "✦ 名下房产数量与估值",
                "✦ 零钱通 / 理财资产规模",
                "✦ 商业保险 & 社保缴纳情况",
                "✦ 家庭负债结构综合评估",
            ],
            "x": 4.75, "y": 2.65, "w": 3.8, "h": 4.2,
        },
        {
            "bg": CARD_BG3, "accent": ORANGE,
            "icon": "📊", "title": "授信敏感度",
            "subtitle": "行为响应 & 风险分层",
            "items": [
                "✦ 历史提额 / 降额响应行为",
                "✦ 用信频率与还款规律性",
                "✦ 逾期预警信号提前识别",
                "✦ 客群精细化敏感度分层",
            ],
            "x": 9.2, "y": 2.65, "w": 3.8, "h": 4.2,
        },
    ]

    for card in cards:
        x, y, w, h = card["x"], card["y"], card["w"], card["h"]
        add_rounded_rect(slide, x, y, w, h,
                         fill_rgb=card["bg"],
                         line_rgb=card["accent"], line_width=Pt(2), corner=0.08)
        # 顶部色条
        add_rect(slide, x, y, w, 0.55, fill_rgb=card["accent"])
        # 圆角遮盖色条顶角——用同色覆盖下半部分的线框（视觉技巧，留线色）
        add_text_box(slide, x + 0.12, y + 0.07, w - 0.2, 0.42,
                     f"{card['icon']}  {card['title']}", size=13, bold=True, color=WHITE)
        add_text_box(slide, x + 0.12, y + 0.62, w - 0.2, 0.35,
                     card["subtitle"], size=10, bold=False,
                     color=card["accent"])
        # 分割线
        add_rect(slide, x + 0.15, y + 1.0, w - 0.3, 0.03, fill_rgb=card["accent"])
        for i, item in enumerate(card["items"]):
            add_text_box(slide, x + 0.15, y + 1.12 + i * 0.72, w - 0.3, 0.62,
                         item, size=10.5, color=MID_TEXT)

        # 连接线（从中央节点引出）
        node_cx = 4.8 + 3.7 / 2  # 节点中心x
        node_by = 1.45 + 0.85      # 节点底部y
        card_top_cx = x + w / 2
        add_connector(slide, node_cx, node_by, card_top_cx, y,
                      color=card["accent"], width=Pt(1.5))

    # 底部说明
    add_rounded_rect(slide, 0.3, 6.88, W - 0.6, 0.5,
                     fill_rgb=BADGE_BG, line_rgb=DIVIDER, line_width=Pt(1), corner=0.06)
    add_text_box(slide, 0.5, 6.92, W - 1.0, 0.38,
                 "三维体系相互印证，综合判断授信额度与风险定价，有效提升工商经营客群的风控精准度与客户留存率。",
                 size=9.5, color=MID_TEXT)

    return slide


# ── 幻灯片 4：企业基本面详解 ──────────────────────────────────────────────────
def make_enterprise(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    W, H = 13.33, 7.5

    add_rect(slide, 0, 0, W, 1.1, fill_rgb=DEEP_BLUE)
    add_rect(slide, 0, 1.1, W, 0.06, fill_rgb=ACCENT_BLUE)
    add_text_box(slide, 0.4, 0.22, W - 0.8, 0.65,
                 "03  企业基本面  —  经营合规性与规模评估", size=20, bold=True, color=WHITE)
    add_text_box(slide, 0.4, 0.72, W - 0.8, 0.32,
                 "新增工商特许经营 · 分支机构 · 对外投资三大维度，构建企业经营画像",
                 size=11, color=RGBColor(0xAE, 0xD6, 0xF1))

    details = [
        {
            "accent": ACCENT_BLUE, "bg": CARD_BG1,
            "no": "01", "title": "工商特许经营资质",
            "tags": ["餐饮/零售/连锁", "资质年限校验", "合规经营加分"],
            "desc": (
                "通过工商数据识别客户是否持有特许经营权，"
                "包括餐饮、零售、连锁门店等品类。\n"
                "持有有效许可证且合规经营期 ≥ 3 年，在评分卡中给予经营稳定性加权。"
            ),
        },
        {
            "accent": RGBColor(0x17, 0x7E, 0xC2), "bg": RGBColor(0xE3, 0xF2, 0xFD),
            "no": "02", "title": "分支机构布局分析",
            "tags": ["机构数量", "地域覆盖度", "扩张趋势"],
            "desc": (
                "统计企业登记在册的分支机构数量与跨城市分布，"
                "量化经营规模指数。\n"
                "近12个月新增分支机构视为正向扩张信号，可适当提升授信上限。"
            ),
        },
        {
            "accent": RGBColor(0x0E, 0x6B, 0xA8), "bg": RGBColor(0xD9, 0xEE, 0xFA),
            "no": "03", "title": "对外投资与关联图谱",
            "tags": ["被投企业资质", "持股比例", "关联风险穿透"],
            "desc": (
                "挖掘客户作为股东的对外投资记录，评估被投企业的经营状况与风险。\n"
                "对存在多层关联、被投企业失信或注销等情形进行负面标记，触发人工复核。"
            ),
        },
    ]

    for idx, d in enumerate(details):
        x, y, w, h = 0.3 + idx * 4.35, 1.38, 4.1, 5.8
        add_rounded_rect(slide, x, y, w, h,
                         fill_rgb=d["bg"], line_rgb=d["accent"],
                         line_width=Pt(1.5), corner=0.08)
        # 顶部序号区
        add_rect(slide, x, y, w, 0.75, fill_rgb=d["accent"])
        add_text_box(slide, x + 0.12, y + 0.06, 0.55, 0.6,
                     d["no"], size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0x99))
        add_text_box(slide, x + 0.72, y + 0.18, w - 0.85, 0.45,
                     d["title"], size=12, bold=True, color=WHITE)
        # 标签
        tag_x = x + 0.15
        for tag in d["tags"]:
            tag_w = len(tag) * 0.165 + 0.2
            add_rounded_rect(slide, tag_x, y + 0.9, tag_w, 0.32,
                             fill_rgb=WHITE, line_rgb=d["accent"],
                             line_width=Pt(1), corner=0.15)
            add_text_box(slide, tag_x + 0.05, y + 0.93, tag_w - 0.05, 0.28,
                         tag, size=8.5, color=d["accent"])
            tag_x += tag_w + 0.12
        # 正文
        add_text_box(slide, x + 0.15, y + 1.38, w - 0.3, 3.9,
                     d["desc"], size=10.5, color=MID_TEXT)

    return slide


# ── 幻灯片 5：个人基本面详解 ──────────────────────────────────────────────────
def make_personal(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    W, H = 13.33, 7.5

    add_rect(slide, 0, 0, W, 1.1, fill_rgb=DEEP_BLUE)
    add_rect(slide, 0, 1.1, W, 0.06, fill_rgb=GREEN)
    add_text_box(slide, 0.4, 0.22, W - 0.8, 0.65,
                 "04  个人基本面  —  资产负债与还款能力", size=20, bold=True, color=WHITE)
    add_text_box(slide, 0.4, 0.72, W - 0.8, 0.32,
                 "综合考量房产 · 零钱通/理财 · 个人保险 · 社保，构建个人偿债能力画像",
                 size=11, color=RGBColor(0xAE, 0xD6, 0xF1))

    items = [
        {
            "icon": "🏠", "title": "房产资产",
            "accent": GREEN, "bg": CARD_BG2,
            "points": [
                "名下住宅/商铺数量与估值",
                "抵押状态与净资产测算",
                "房产位于核心城市给予加权",
                "共有产权需核实实际控制权",
            ],
        },
        {
            "icon": "💰", "title": "零钱通 / 理财",
            "accent": RGBColor(0x27, 0xAE, 0x60), "bg": RGBColor(0xE2, 0xF9, 0xED),
            "points": [
                "近6个月平均余额作为流动性指标",
                "定期理财规模反映风险偏好",
                "资产波动率监控异常资金流入",
                "余额持续下降触发预警标记",
            ],
        },
        {
            "icon": "🛡", "title": "个人保险",
            "accent": RGBColor(0x1E, 0x8B, 0x72), "bg": RGBColor(0xD5, 0xF5, 0xEA),
            "points": [
                "寿险 / 意外险保额与保单状态",
                "年缴保费反映稳定收入证明",
                "保险公司评级筛选有效保单",
                "保单质押情况纳入负债计算",
            ],
        },
        {
            "icon": "📋", "title": "社会保险",
            "accent": RGBColor(0x16, 0x73, 0x5E), "bg": RGBColor(0xC8, 0xF0, 0xE1),
            "points": [
                "连续缴纳月数与缴纳基数",
                "中断记录识别经营断层风险",
                "多城市缴纳记录识别经营扩张",
                "与工商注册城市一致性验证",
            ],
        },
    ]

    for idx, item in enumerate(items):
        x, y, w, h = 0.3 + idx * 3.2, 1.38, 2.95, 5.8
        add_rounded_rect(slide, x, y, w, h,
                         fill_rgb=item["bg"], line_rgb=item["accent"],
                         line_width=Pt(1.5), corner=0.08)
        add_rect(slide, x, y, w, 0.7, fill_rgb=item["accent"])
        add_text_box(slide, x + 0.12, y + 0.08, 0.5, 0.55,
                     item["icon"], size=18)
        add_text_box(slide, x + 0.65, y + 0.16, w - 0.75, 0.42,
                     item["title"], size=12, bold=True, color=WHITE)
        for i, pt in enumerate(item["points"]):
            py = y + 0.88 + i * 1.2
            add_rounded_rect(slide, x + 0.12, py, w - 0.24, 1.05,
                             fill_rgb=WHITE, line_rgb=DIVIDER, line_width=Pt(0.75), corner=0.06)
            add_text_box(slide, x + 0.22, py + 0.1, w - 0.44, 0.82,
                         pt, size=9.5, color=MID_TEXT)

    return slide


# ── 幻灯片 6：授信敏感度详解 ──────────────────────────────────────────────────
def make_sensitivity(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    W, H = 13.33, 7.5

    add_rect(slide, 0, 0, W, 1.1, fill_rgb=DEEP_BLUE)
    add_rect(slide, 0, 1.1, W, 0.06, fill_rgb=ORANGE)
    add_text_box(slide, 0.4, 0.22, W - 0.8, 0.65,
                 "05  授信敏感度  —  行为响应与风险分层", size=20, bold=True, color=WHITE)
    add_text_box(slide, 0.4, 0.72, W - 0.8, 0.32,
                 "依据客户对授信变动的行为响应模式，建立差异化授信策略",
                 size=11, color=RGBColor(0xAE, 0xD6, 0xF1))

    # 左侧：分层说明
    add_rounded_rect(slide, 0.3, 1.38, 5.5, 5.8,
                     fill_rgb=CARD_BG3, line_rgb=ORANGE, line_width=Pt(1.5), corner=0.08)
    add_rect(slide, 0.3, 1.38, 5.5, 0.55, fill_rgb=ORANGE)
    add_text_box(slide, 0.5, 1.44, 5.1, 0.42,
                 "敏感度分层逻辑", size=12, bold=True, color=WHITE)

    tiers = [
        ("高敏感客户", ORANGE,
         "• 授信提升 → 立即大额支取\n• 降额 → 立即申请他行额度\n"
         "• 策略：保持额度稳定，避免触发流失"),
        ("中敏感客户", RGBColor(0xCA, 0x6F, 0x1E),
         "• 对额度变化反应适中\n• 还款规律，偶有延期\n"
         "• 策略：阶梯式提额，绑定更多产品"),
        ("低敏感客户", RGBColor(0x9A, 0x7D, 0x0A),
         "• 额度变化对行为影响小\n• 高忠诚度，长期活跃\n"
         "• 策略：挖掘交叉销售，提升综合贡献"),
    ]
    for idx, (name, color, desc) in enumerate(tiers):
        y = 2.12 + idx * 1.7
        add_rounded_rect(slide, 0.5, y, 5.0, 1.5,
                         fill_rgb=WHITE, line_rgb=color, line_width=Pt(1.5), corner=0.06)
        add_rect(slide, 0.5, y, 0.08, 1.5, fill_rgb=color)
        add_text_box(slide, 0.7, y + 0.1, 4.6, 0.38,
                     name, size=11, bold=True, color=color)
        add_text_box(slide, 0.7, y + 0.5, 4.6, 0.88,
                     desc, size=9.5, color=MID_TEXT)

    # 右侧：指标体系
    add_rounded_rect(slide, 6.1, 1.38, 6.9, 5.8,
                     fill_rgb=WHITE, line_rgb=DIVIDER, line_width=Pt(1), corner=0.08)
    add_rect(slide, 6.1, 1.38, 6.9, 0.55, fill_rgb=MID_BLUE)
    add_text_box(slide, 6.3, 1.44, 6.5, 0.42,
                 "核心监测指标", size=12, bold=True, color=WHITE)

    metrics = [
        ("提额响应率", "提额后30天内支取金额 / 新增额度",
         "≥ 80%  →  高敏感", ACCENT_BLUE),
        ("还款规律性得分", "连续按时还款期数 / 观察期总期数",
         "≥ 90%  →  优质信号", GREEN),
        ("用信频率指数", "月均用信次数 × 单笔平均金额",
         "上升趋势 → 经营扩张", ORANGE),
        ("逾期预警信号", "近3期最小还款日提前天数均值",
         "< 3天  →  触发预警", RGBColor(0xC0, 0x39, 0x2B)),
    ]
    for idx, (name, formula, benchmark, color) in enumerate(metrics):
        y = 2.1 + idx * 1.42
        add_rounded_rect(slide, 6.3, y, 6.5, 1.25,
                         fill_rgb=BADGE_BG, line_rgb=color, line_width=Pt(1), corner=0.06)
        add_text_box(slide, 6.5, y + 0.08, 6.1, 0.35,
                     name, size=11, bold=True, color=color)
        add_text_box(slide, 6.5, y + 0.42, 6.1, 0.32,
                     f"计算：{formula}", size=9, color=MID_TEXT)
        add_text_box(slide, 6.5, y + 0.72, 6.1, 0.32,
                     f"基准：{benchmark}", size=9, color=MID_TEXT,)

    return slide


# ── 幻灯片 7：实施路径 ────────────────────────────────────────────────────────
def make_roadmap(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    W, H = 13.33, 7.5

    add_rect(slide, 0, 0, W, 1.1, fill_rgb=DEEP_BLUE)
    add_rect(slide, 0, 1.1, W, 0.06, fill_rgb=ORANGE)
    add_text_box(slide, 0.4, 0.22, W - 0.8, 0.65,
                 "06  实施路径与预期收益", size=22, bold=True, color=WHITE)
    add_text_box(slide, 0.4, 0.72, W - 0.8, 0.32,
                 "三阶段推进，逐步构建精准授信体系", size=11,
                 color=RGBColor(0xAE, 0xD6, 0xF1))

    phases = [
        {
            "no": "Phase 1", "period": "第 1-2 月",
            "title": "数据接入与画像构建",
            "accent": ACCENT_BLUE, "bg": CARD_BG1,
            "tasks": [
                "对接工商数据库 API",
                "建立企业经营评分卡 v1",
                "社保 / 保险数据清洗入库",
                "完成客群基础画像建模",
            ],
            "kpi": "覆盖率 ≥ 85%",
        },
        {
            "no": "Phase 2", "period": "第 3-4 月",
            "title": "敏感度建模与策略联调",
            "accent": GREEN, "bg": CARD_BG2,
            "tasks": [
                "训练授信敏感度分类模型",
                "分层策略规则配置上线",
                "A/B 测试验证策略效果",
                "风险监控看板搭建",
            ],
            "kpi": "提额审批效率 +30%",
        },
        {
            "no": "Phase 3", "period": "第 5-6 月",
            "title": "全量上线与持续优化",
            "accent": ORANGE, "bg": CARD_BG3,
            "tasks": [
                "全量客群切换新评分体系",
                "逾期预警模型闭环验证",
                "季度模型迭代机制建立",
                "不良率与客户留存监控",
            ],
            "kpi": "不良率目标 ≤ 3.5%",
        },
    ]

    for idx, ph in enumerate(phases):
        x, y, w, h = 0.3 + idx * 4.35, 1.38, 4.1, 5.0
        add_rounded_rect(slide, x, y, w, h,
                         fill_rgb=ph["bg"], line_rgb=ph["accent"],
                         line_width=Pt(2), corner=0.08)
        add_rect(slide, x, y, w, 0.78, fill_rgb=ph["accent"])
        add_text_box(slide, x + 0.12, y + 0.05, w - 0.2, 0.35,
                     ph["no"], size=10, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0x99))
        add_text_box(slide, x + 0.12, y + 0.36, w - 0.2, 0.32,
                     f"（{ph['period']}）{ph['title']}", size=10.5, bold=True, color=WHITE)
        for i, task in enumerate(ph["tasks"]):
            add_rounded_rect(slide, x + 0.15, y + 0.95 + i * 0.95, w - 0.3, 0.78,
                             fill_rgb=WHITE, line_rgb=DIVIDER,
                             line_width=Pt(0.75), corner=0.05)
            add_text_box(slide, x + 0.28, y + 1.05 + i * 0.95, w - 0.55, 0.55,
                         task, size=10, color=MID_TEXT)

        # KPI 徽章
        add_rounded_rect(slide, x + 0.15, y + 4.58, w - 0.3, 0.55,
                         fill_rgb=ph["accent"], line_rgb=ph["accent"],
                         line_width=Pt(0), corner=0.1)
        add_text_box(slide, x + 0.2, y + 4.67, w - 0.4, 0.38,
                     f"🎯  目标：{ph['kpi']}", size=10, bold=True, color=WHITE)

    # 箭头（文字替代）
    add_text_box(slide, 4.43, 3.52, 0.45, 0.4,
                 "➜", size=22, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 8.78, 3.52, 0.45, 0.4,
                 "➜", size=22, color=GREEN, align=PP_ALIGN.CENTER)

    return slide


# ── 主函数 ────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    make_cover(prs)
    make_background(prs)
    make_framework(prs)
    make_enterprise(prs)
    make_personal(prs)
    make_sensitivity(prs)
    make_roadmap(prs)

    out = "/tmp/credit_framework.pptx"
    prs.save(out)
    print(f"✅  已生成：{out}")
    print(f"   共 {len(prs.slides)} 张幻灯片")
    print("   封面 → 客群背景 → 整体框架 → 企业基本面 → 个人基本面 → 授信敏感度 → 实施路径")


if __name__ == "__main__":
    main()

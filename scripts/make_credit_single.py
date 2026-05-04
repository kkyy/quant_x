"""
单页授信框架图示（美化版）
原图逻辑：客群特点 → 应对方案（三维体系）
输出：/tmp/credit_single.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── 颜色 ──────────────────────────────────────────────────────────────────────
DEEP_BLUE   = RGBColor(0x1A, 0x3A, 0x6B)
MID_BLUE    = RGBColor(0x25, 0x5F, 0xAB)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE  = RGBColor(0xD6, 0xEA, 0xF8)
ORANGE      = RGBColor(0xE8, 0x7E, 0x04)
GREEN       = RGBColor(0x1A, 0x8A, 0x5A)
AMBER       = RGBColor(0xF3, 0x9C, 0x12)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT   = RGBColor(0x55, 0x55, 0x66)
PALE_BG     = RGBColor(0xF4, 0xF7, 0xFD)
DIVIDER     = RGBColor(0xCC, 0xD9, 0xEA)


def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill:
        s.fill.fore_color.rgb = fill
    s.line.fill.background() if not line else None
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    s.shadow.inherit = False
    return s


def rrect(slide, l, t, w, h, fill=None, line=None, lw=Pt(1), corner=0.08):
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    prstGeom = s.element.spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is not None:
            for gd in avLst.findall(qn('a:gd')):
                if gd.get('name') == 'adj':
                    gd.set('fmla', f'val {int(corner * 50000)}')
    s.fill.solid() if fill else s.fill.background()
    if fill:
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def tb(slide, l, t, w, h, text, size=11, bold=False, color=None,
       align=PP_ALIGN.LEFT, wrap=True):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color else RGBColor(0x1C, 0x1C, 0x1C)
    return tx


def connector(slide, x1, y1, x2, y2, color=ACCENT_BLUE, w=Pt(1.5)):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = w
    return c


# ── 布局常量 ──────────────────────────────────────────────────────────────────
W, H = 13.33, 7.5


def make_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── 背景 ──────────────────────────────────────────────────────────────────
    rect(slide, 0, 0, W, H, fill=PALE_BG)

    # ── 顶部标题栏 ────────────────────────────────────────────────────────────
    rect(slide, 0, 0, W, 0.95, fill=DEEP_BLUE)
    rect(slide, 0, 0.95, W, 0.055, fill=ORANGE)
    tb(slide, 0.4, 0.14, 9, 0.5,
       "工商经营客群  ·  授信框架", size=20, bold=True, color=WHITE)
    tb(slide, 0.4, 0.62, 9, 0.28,
       "强化经营特质挖掘，引入授信敏感度调整", size=10.5, color=RGBColor(0xAE, 0xD6, 0xF1))

    # ── 分区标注（右上） ──────────────────────────────────────────────────────
    tb(slide, 10.8, 0.2, 2.3, 0.28,
       "客群框架  /  Framework", size=9, color=RGBColor(0x88, 0xAB, 0xD5),
       align=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════════════════════════════════════
    # 左侧：客群特点（背景 & 痛点）
    # ════════════════════════════════════════════════════════════════════════
    LEFT_X, CONTENT_Y, CARD_H = 0.3, 1.18, 5.9

    rrect(slide, LEFT_X, CONTENT_Y, 3.6, CARD_H,
          fill=WHITE, line=ACCENT_BLUE, lw=Pt(1.5), corner=0.06)
    # 顶色条
    rect(slide, LEFT_X, CONTENT_Y, 3.6, 0.55, fill=MID_BLUE)
    tb(slide, LEFT_X + 0.15, CONTENT_Y + 0.1, 3.3, 0.38,
       "客群特点", size=12, bold=True, color=WHITE)

    pain_items = [
        (ORANGE,      "无强收入数据源",
         "缺乏稳定工资流水，\n传统收入核验方式失效"),
        (RGBColor(0xC0, 0x39, 0x2B), "风险偏高",
         "年化不良率约 4.3%，\n高于普通零售贷款"),
        (ACCENT_BLUE, "工商经营属性强",
         "有明确工商注册记录，\n经营维度信息可挖掘"),
    ]
    for idx, (dot_color, title, desc) in enumerate(pain_items):
        py = CONTENT_Y + 0.72 + idx * 1.72
        rrect(slide, LEFT_X + 0.18, py, 3.1, 1.5,
              fill=PALE_BG, line=DIVIDER, lw=Pt(0.75), corner=0.06)
        # 色点
        rrect(slide, LEFT_X + 0.28, py + 0.14, 0.22, 0.22,
              fill=dot_color, corner=0.5)
        tb(slide, LEFT_X + 0.56, py + 0.1, 2.6, 0.3,
           title, size=10.5, bold=True, color=dot_color)
        tb(slide, LEFT_X + 0.28, py + 0.44, 2.85, 0.9,
           desc, size=9.5, color=GRAY_TEXT)

    # ════════════════════════════════════════════════════════════════════════
    # 中间：过渡箭头 + 核心命题
    # ════════════════════════════════════════════════════════════════════════
    ARROW_X = 4.15
    # 竖线
    connector(slide, ARROW_X + 0.2, CONTENT_Y + 0.6,
              ARROW_X + 0.2, CONTENT_Y + CARD_H - 0.6,
              color=DIVIDER, w=Pt(1))
    # 箭头文字
    tb(slide, ARROW_X - 0.05, CONTENT_Y + CARD_H / 2 - 0.5, 0.52, 1.0,
       "⟹", size=26, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(slide, ARROW_X - 0.12, CONTENT_Y + CARD_H / 2 + 0.35, 0.75, 0.5,
       "应对\n方案", size=8, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # 右侧：三维授信体系
    # ════════════════════════════════════════════════════════════════════════
    RIGHT_X = 4.72

    # 外框
    rrect(slide, RIGHT_X, CONTENT_Y, 8.3, CARD_H,
          fill=WHITE, line=ACCENT_BLUE, lw=Pt(1.5), corner=0.06)
    rect(slide, RIGHT_X, CONTENT_Y, 8.3, 0.55, fill=MID_BLUE)
    tb(slide, RIGHT_X + 0.18, CONTENT_Y + 0.1, 7.9, 0.38,
       "授信评估体系  —  三维框架", size=12, bold=True, color=WHITE)

    # 三维卡片配置
    CARD_CONFIGS = [
        {
            "accent": ACCENT_BLUE,
            "bg": RGBColor(0xEB, 0xF5, 0xFB),
            "icon": "🏢",
            "title": "企业基本面",
            "subtitle": "经营合规性 & 规模评估",
            "items": [
                "新增工商特许经营资质核验",
                "分支机构数量与地域覆盖",
                "对外投资及关联企业图谱",
            ],
        },
        {
            "accent": GREEN,
            "bg": RGBColor(0xE8, 0xF8, 0xF0),
            "icon": "👤",
            "title": "个人基本面",
            "subtitle": "资产负债 & 还款能力",
            "items": [
                "综合考量名下房产及估值",
                "零钱通 / 理财资产规模",
                "个人保险 & 社保缴纳情况",
            ],
        },
        {
            "accent": ORANGE,
            "bg": RGBColor(0xFE, 0xF9, 0xE7),
            "icon": "📊",
            "title": "授信敏感度",
            "subtitle": "行为响应 & 风险分层",
            "items": [
                "历史提额 / 降额响应行为",
                "用信频率与还款规律性",
                "逾期预警信号精细化分层",
            ],
        },
    ]

    CARD_W = 2.45
    CARD_Y_START = CONTENT_Y + 0.72
    CARD_H_INNER = 4.98

    for idx, cfg in enumerate(CARD_CONFIGS):
        cx = RIGHT_X + 0.22 + idx * (CARD_W + 0.27)
        rrect(slide, cx, CARD_Y_START, CARD_W, CARD_H_INNER,
              fill=cfg["bg"], line=cfg["accent"], lw=Pt(1.5), corner=0.07)
        # 顶色条
        rect(slide, cx, CARD_Y_START, CARD_W, 0.62, fill=cfg["accent"])
        # Icon + title
        tb(slide, cx + 0.1, CARD_Y_START + 0.06, 0.45, 0.5,
           cfg["icon"], size=16, align=PP_ALIGN.CENTER)
        tb(slide, cx + 0.55, CARD_Y_START + 0.1, CARD_W - 0.62, 0.3,
           cfg["title"], size=11, bold=True, color=WHITE)
        tb(slide, cx + 0.55, CARD_Y_START + 0.38, CARD_W - 0.62, 0.22,
           cfg["subtitle"], size=8, color=RGBColor(0xD6, 0xEA, 0xF8))
        # 分割线
        rect(slide, cx + 0.12, CARD_Y_START + 0.74, CARD_W - 0.24, 0.028,
             fill=cfg["accent"])
        # 条目
        for i, item in enumerate(cfg["items"]):
            iy = CARD_Y_START + 0.86 + i * 1.32
            # 小色圆
            rrect(slide, cx + 0.14, iy + 0.08, 0.16, 0.16,
                  fill=cfg["accent"], corner=0.5)
            tb(slide, cx + 0.36, iy, CARD_W - 0.45, 0.95,
               item, size=9.5, color=GRAY_TEXT)

    # ── 底部说明条 ────────────────────────────────────────────────────────────
    rrect(slide, 0.3, 7.18, W - 0.6, 0.24,
          fill=LIGHT_BLUE, line=DIVIDER, lw=Pt(0.75), corner=0.04)
    tb(slide, 0.5, 7.2, W - 1.0, 0.2,
       "三维体系相互印证，综合判断授信额度与风险定价，有效提升风控精准度与客户留存率",
       size=8.5, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    return slide


def main():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    make_slide(prs)
    out = "/tmp/credit_single.pptx"
    prs.save(out)
    print(f"✅ 已生成：{out}")


if __name__ == "__main__":
    main()
